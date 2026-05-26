from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

OUT = Path('/Users/tachongrak/Projects/hr/docs/presentation/humi-system-design-sample.pptx')
OUT.parent.mkdir(parents=True, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

C = {
    'primary': '1FA8A0', 'primarySoft': 'D6EEEC', 'primaryHover': '126E69',
    'secondary': '5B6CE0', 'secondarySoft': 'E1E4FB',
    'canvas': 'F6F1E8', 'canvasSoft': 'FCFAF5', 'surface': 'FFFFFF',
    'ink': '0E1B2C', 'inkSoft': '243447', 'inkMuted': '5A6A7E', 'hairline': 'E7DFD1',
    'butter': 'E8C46B', 'successSoft': 'D1FAE5', 'warningSoft': 'FEF3C7',
    'dangerSoft': 'FFEDD5', 'dangerInk': '9A3412', 'warmText': 'E7E3D8'
}
FONT = 'Aptos'
FONT_HEAD = 'Aptos Display'

def rgb(h):
    h = h.replace('#','')
    return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:6],16))

def set_fill(shape, color, trans=0):
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(C.get(color, color)); shape.fill.transparency = trans

def no_line(shape):
    shape.line.fill.background()

def set_line(shape, color='hairline', width=1):
    shape.line.color.rgb = rgb(C.get(color, color)); shape.line.width = Pt(width)

def bg(slide, color='canvas'):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_fill(sh, color); no_line(sh)
    return sh

def text(slide, value, x, y, w, h, size=14, color='ink', bold=False, align=PP_ALIGN.LEFT, font=FONT, margin=0.04, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(margin); tf.margin_right = Inches(margin); tf.margin_top = Inches(margin); tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = value
    run.font.name = font; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = rgb(C.get(color, color))
    return box

def rect(slide, x, y, w, h, fill='surface', radius=True, border=None, trans=0):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(sh, fill, trans)
    if border: set_line(sh, border, 0.8)
    else: no_line(sh)
    return sh

def pill(slide, label, x, y, w, fill='primarySoft', color='ink'):
    rect(slide, x, y, w, 0.34, fill, True)
    text(slide, label, x+0.05, y+0.085, w-0.1, 0.1, 8.5, color, True, PP_ALIGN.CENTER, margin=0)

def title(slide, t, st=None, dark=False):
    text(slide, t, 0.65, 0.42, 8.8, 0.55, 27, 'surface' if dark else 'ink', True, font=FONT_HEAD, margin=0)
    if st:
        text(slide, st, 0.67, 0.98, 8.5, 0.26, 10.5, 'warmText' if dark else 'inkMuted', margin=0)

def footer(slide, n, dark=False):
    text(slide, f'Humi System Design · DESIGN.md · {n:02d}', 0.65, 7.08, 4.2, 0.16, 7.5, 'warmText' if dark else 'inkMuted', margin=0)

# Slide 1 cover
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, 'ink')
for x,y,w,c,tr in [(9.2,-0.55,4.4,'primary',0),(10.7,3.85,2.4,'butter',8),(8.25,4.9,1.25,'secondary',0)]:
    o=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(w)); set_fill(o,c,tr); no_line(o)
pill(s, 'DESIGN SYSTEM', 0.78, 0.7, 1.72, 'primary', 'ink')
text(s, 'Humi System Design', 0.78, 1.55, 6.8, 0.65, 36, 'surface', True, font=FONT_HEAD, margin=0)
text(s, 'Cream canvas + navy ink + teal action system\nfor UI, mockups, and executive presentations', 0.82, 2.45, 6.25, 0.58, 16, 'warmText', margin=0)
for i,(num,lab,col) in enumerate([('17','curated colors','primary'),('18','usable components','butter'),('0','lint warnings','secondary')]):
    x=0.82+i*2.0
    rect(s,x,4.55,1.72,0.9,'inkSoft',True,'inkSoft')
    text(s,num,x+0.16,4.68,0.48,0.26,21,col,True,font=FONT_HEAD,margin=0)
    text(s,lab,x+0.70,4.72,0.83,0.18,8.5,'warmText',margin=0)
rect(s,8.48,1.22,3.65,3.95,'surface',True)
text(s,'Presentation rhythm',8.82,1.55,2.2,0.28,17,'ink',True,margin=0)
for j,(a,b,c) in enumerate([('Ink cover','executive framing','ink'),('Cream content','calm readable canvas','canvas'),('White cards','modular UI blocks','surface'),('Teal decisions','primary action signal','primary'),('Butter KPIs','warm emphasis','butter')]):
    y=2.02+j*0.55
    rect(s,8.85,y,0.34,0.34,c,True,'hairline')
    text(s,a,9.35,y+0.01,1.15,0.13,9.2,'ink',True,margin=0)
    text(s,b,10.40,y+0.01,1.35,0.13,8,'inkMuted',margin=0)
footer(s,1,True)

# Slide 2 token kit
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'canvas'); title(s,'Curated token kit','คัดเฉพาะ palette และ components ที่ใช้งานจริงใน UI / presentation'); footer(s,2)
palette=[('primary','#1FA8A0','CTA / active'),('secondary','#5B6CE0','info / flow'),('canvas','#F6F1E8','page bg'),('surface','#FFFFFF','cards'),('ink','#0E1B2C','text / sidebar'),('butter','#E8C46B','KPI'),('dangerSoft','#FFEDD5','safe danger'),('hairline','#E7DFD1','divider')]
for i,(name,hexv,use) in enumerate(palette):
    x=0.78+(i%4)*3.05; y=1.55+(i//4)*1.28
    rect(s,x,y,2.55,0.92,'surface',True,'hairline')
    rect(s,x+0.18,y+0.17,0.58,0.58,name,True,'hairline')
    text(s,name,x+0.92,y+0.18,1.0,0.16,10.5,'ink',True,margin=0)
    text(s,hexv,x+0.92,y+0.43,0.85,0.14,8.5,'inkMuted',margin=0)
    text(s,use,x+1.66,y+0.33,0.68,0.16,8.2,'inkMuted',False,PP_ALIGN.RIGHT,margin=0)
rect(s,0.78,4.55,11.72,1.25,'ink',True)
text(s,'Rule of thumb',1.08,4.88,1.65,0.18,13,'butter',True,margin=0)
text(s,'ใช้สีเท่าที่จำเป็น: teal = decision/action, indigo = flow/info, butter = KPI, pumpkin soft = danger without red.',2.75,4.85,8.95,0.22,13.2,'surface',True,margin=0)

# Slide 3 UI grammar with dashboard mock
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'canvas'); title(s,'UI grammar: Humi page anatomy','Cream canvas, dark rail, white cards, soft status tags'); footer(s,3)
# mock app frame
rect(s,0.8,1.4,11.8,4.75,'surface',True,'hairline')
rect(s,0.8,1.4,2.15,4.75,'ink',False)
text(s,'Humi',1.05,1.72,0.95,0.22,16,'surface',True,font=FONT_HEAD,margin=0)
for i,lab in enumerate(['Employee','Time','Payroll','Benefits','Reports']):
    y=2.22+i*0.48
    rect(s,1.02,y,1.55,0.30,'primary' if i==1 else 'inkSoft',True)
    text(s,lab,1.20,y+0.08,1.0,0.08,7.8,'ink' if i==1 else 'warmText',True if i==1 else False,margin=0)
text(s,'Time Dashboard',3.28,1.75,2.5,0.24,18,'ink',True,margin=0)
text(s,'Today · HR sign-off view',3.30,2.07,2.0,0.14,9,'inkMuted',margin=0)
for i,(n,l,c) in enumerate([('92%','Timesheet complete','primarySoft'),('14','Pending approval','warningSoft'),('3','Returned items','dangerSoft')]):
    x=3.28+i*2.1
    rect(s,x,2.52,1.75,0.88,c,True)
    text(s,n,x+0.15,2.70,0.6,0.22,20,'ink',True,font=FONT_HEAD,margin=0)
    text(s,l,x+0.78,2.74,0.75,0.16,7.8,'inkSoft',margin=0)
rect(s,3.28,3.72,5.65,1.75,'canvasSoft',True,'hairline')
for i,(emp,stat,c) in enumerate([('Somchai P.','Approved','successSoft'),('Narin K.','Pending','warningSoft'),('Mali S.','Return','dangerSoft')]):
    y=3.97+i*0.45
    text(s,emp,3.55,y,1.1,0.12,8.5,'ink',margin=0)
    rect(s,5.75,y-0.04,0.95,0.22,c,True)
    text(s,stat,5.83,y+0.02,0.68,0.08,6.6,'dangerInk' if c=='dangerSoft' else 'ink',False,PP_ALIGN.CENTER,margin=0)
    set_line(s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.48), Inches(y+0.28), Inches(8.6), Inches(y+0.28)), 'hairline', 0.5)
rect(s,9.35,2.52,2.55,2.95,'primarySoft',True)
text(s,'Presenter note',9.66,2.82,1.45,0.18,11,'ink',True,margin=0)
text(s,'Show component hierarchy first, then business status. Avoid raw hex and legacy red.',9.66,3.18,1.65,0.72,9.5,'inkSoft',margin=0)

# Slide 4 presentation recipe
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'canvas'); title(s,'Presentation recipe','ใช้ Humi tokens ทำ executive deck ให้ดูเหมือน product system เดียวกัน'); footer(s,4)
steps=[('1','Cover','Ink background + teal/butter accents','ink'),('2','Decision','Cream page + white cards + teal CTA','primary'),('3','Flow','Indigo nodes + hairline connectors','secondary'),('4','KPI','Butter metric cards, mono numbers','butter'),('5','Risks','Danger-soft tags, no red','dangerSoft')]
for i,(num,h,b,c) in enumerate(steps):
    x=0.8+i*2.42
    circ=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.62),Inches(1.75),Inches(0.72),Inches(0.72)); set_fill(circ,c); set_line(circ,'hairline',0.5)
    text(s,num,x+0.62,1.93,0.72,0.1,12,'surface' if c in ['ink','secondary'] else 'ink',True,PP_ALIGN.CENTER,margin=0)
    if i<4: set_line(s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x+1.42), Inches(2.11), Inches(x+2.25), Inches(2.11)), 'hairline', 1.4)
    text(s,h,x+0.25,2.72,1.45,0.18,12,'ink',True,PP_ALIGN.CENTER,margin=0)
    text(s,b,x+0.10,3.03,1.76,0.42,9.4,'inkMuted',False,PP_ALIGN.CENTER,margin=0)
rect(s,1.1,4.68,10.95,1.18,'surface',True,'hairline')
text(s,'Design prompt for agents',1.42,4.98,1.78,0.16,10.5,'primary',True,margin=0)
text(s,'“Use Humi System Design: cream canvas, navy ink, teal decisions, indigo flow, butter KPIs, rounded white cards, no red family.”',3.18,4.92,7.95,0.24,13.5,'ink',True,margin=0)

# Slide 5 applied executive slide
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'canvas'); title(s,'Applied example: executive-ready slide','ตัวอย่าง slide ที่ใช้ Humi token แล้วอ่านง่ายสำหรับผู้บริหาร'); footer(s,5)
rect(s,0.82,1.35,4.1,4.75,'ink',True)
pill(s,'EXECUTIVE DECISION',1.18,1.72,1.85,'primary','ink')
text(s,'AI-Native Hiring Framework',1.18,2.35,2.8,0.68,25,'surface',True,font=FONT_HEAD,margin=0)
text(s,'เลือกคนที่คิดเป็น คุม AI ได้ และส่งมอบ value ได้จริง',1.20,3.30,2.95,0.36,13,'warmText',margin=0)
for i,(n,l,c) in enumerate([('3','target hires','primary'),('20','score gate','butter')]):
    x=1.22+i*1.55
    rect(s,x,4.35,1.25,0.72,'inkSoft',True)
    text(s,n,x+0.13,4.50,0.35,0.20,17,c,True,font=FONT_HEAD,margin=0)
    text(s,l,x+0.53,4.54,0.55,0.12,7.3,'warmText',margin=0)
# right content cards
for i,(h,b,c) in enumerate([('Think first','ถามหาปัญหาและ user pain ก่อนเลือก tool','primarySoft'),('Control AI','ใช้ AI เป็น co-pilot แต่ตรวจและรับผิดชอบเอง','secondarySoft'),('Deliver value','เลือก MVP, trade-off และวัด business impact','warningSoft')]):
    y=1.55+i*1.25
    rect(s,5.35,y,6.72,0.95,'surface',True,'hairline')
    rect(s,5.65,y+0.24,0.44,0.44,c,True)
    text(s,h,6.35,y+0.23,1.55,0.18,13,'ink',True,margin=0)
    text(s,b,7.82,y+0.23,3.55,0.22,10.2,'inkSoft',margin=0)
rect(s,5.35,5.48,6.72,0.58,'primary',True)
text(s,'Humi presentation style = fewer words, stronger hierarchy, token-consistent visuals',5.68,5.67,6.0,0.1,10.5,'ink',True,PP_ALIGN.CENTER,margin=0)

# Slide 6 adoption checklist
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'ink'); title(s,'Use it as the system prompt for UI and decks','DESIGN.md is now clean enough to inject into agents or presentation generators',dark=True); footer(s,6,True)
for i,(h,b,c) in enumerate([('UI generation','Start from Humi primitives: page, sidebar, card, input, tags','primary'),('Presentation','Use cover/content/flow/KPI recipe and avoid dense text','butter'),('QA guardrail','Lint DESIGN.md + visual QA contact sheet before delivery','secondary')]):
    x=0.9+i*4.0
    rect(s,x,2.0,3.25,2.7,'surface',True)
    rect(s,x+0.25,2.28,0.54,0.54,c,True)
    text(s,h,x+0.96,2.30,1.55,0.18,14,'ink',True,margin=0)
    text(s,b,x+0.34,3.05,2.45,0.62,11,'inkSoft',margin=0)
rect(s,1.0,5.55,11.2,0.62,'primary',True)
text(s,'Ready: curated DESIGN.md → Tailwind/DTCG export → editable PowerPoint sample → visual QA',1.25,5.75,10.65,0.12,12.5,'ink',True,PP_ALIGN.CENTER,margin=0)

prs.save(OUT)
print(OUT)
