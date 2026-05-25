from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pathlib import Path
import re, os, zipfile

OUT = Path('/Users/tachongrak/Projects/hr/deliverables/hr-app-demo-showcase-2026-05-18.pptx')
OUT.parent.mkdir(parents=True, exist_ok=True)

imgs = {
    'profile': '/Users/tachongrak/Projects/hr/src/frontend/e2e/screenshots/persona-walkthrough/employee-02-profile.png',
    'time': '/Users/tachongrak/Projects/hr/test-artifacts/ultraqa-time-20260518/time-employee-1440.png',
    'payroll': '/Users/tachongrak/Projects/hr/tests/results/screenshots/payroll_setup_verification.png',
    'benefits': '/Users/tachongrak/Projects/hr/src/frontend/e2e/screenshots/persona-walkthrough/employee-03-benefits-hub.png',
    'quick': '/Users/tachongrak/Projects/hr/test-pages/20-quick-approve.png',
    'hrbp': '/Users/tachongrak/Projects/hr/src/frontend/e2e/screenshots/persona-walkthrough/hrbp-01-quick-approve.png',
}

W, H = 13.333333, 7.5
prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]

NAVY = RGBColor(12, 31, 52)
DEEP = RGBColor(8, 22, 38)
TEAL = RGBColor(12, 143, 130)
MINT = RGBColor(105, 211, 194)
AMBER = RGBColor(245, 170, 66)
CREAM = RGBColor(246, 248, 247)
INK = RGBColor(28, 45, 58)
MUTED = RGBColor(96, 112, 124)
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(219, 230, 226)
GREEN = RGBColor(59, 154, 120)


def set_fill(shape, color, transparency=0):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.transparency = transparency
    if hasattr(shape, 'line'):
        shape.line.color.rgb = color


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False, align='left', font='Aptos', transparency=None):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.clear(); tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02); tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
    p.font.name = font; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    return tx


def add_round_rect(slide, x, y, w, h, color=WHITE, line=BORDER, radius=True, transparency=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.fill.transparency = transparency
    shp.line.color.rgb = line; shp.line.width = Pt(1)
    return shp


def add_chip(slide, text, x, y, w, color=TEAL, text_color=WHITE):
    shp = add_round_rect(slide, x, y, w, 0.34, color=color, line=color)
    add_text(slide, text, x+0.12, y+0.065, w-0.24, 0.2, 9.5, text_color, True, 'center')
    return shp


def add_image_card(slide, path, x, y, w, h, label=None, crop_top=0, crop_bottom=0, crop_left=0, crop_right=0):
    add_round_rect(slide, x+0.06, y+0.08, w, h, color=RGBColor(206, 218, 216), line=RGBColor(206, 218, 216), transparency=45)
    frame = add_round_rect(slide, x, y, w, h, color=WHITE, line=BORDER)
    pic = slide.shapes.add_picture(path, Inches(x+0.06), Inches(y+0.06), width=Inches(w-0.12), height=Inches(h-0.12))
    pic.crop_top = crop_top
    pic.crop_bottom = crop_bottom
    pic.crop_left = crop_left
    pic.crop_right = crop_right
    if label:
        add_chip(slide, label, x+0.18, y+0.18, min(2.25, w-0.36), NAVY, WHITE)
    return pic


def add_title(slide, title, subtitle=None, dark=False):
    c = WHITE if dark else INK
    add_text(slide, title, 0.65, 0.45, 7.2, 0.55, 28, c, True)
    if subtitle:
        add_text(slide, subtitle, 0.68, 1.02, 7.8, 0.34, 12.5, MINT if dark else MUTED)


def add_footer(slide, n):
    add_text(slide, f'HR App Demo Showcase  •  {n}/7', 10.5, 7.05, 2.1, 0.22, 8.5, RGBColor(130,145,152), align='right')

# Slide 1 cover
s = prs.slides.add_slide(blank)
set_fill(s.background, DEEP)
# dark panel
add_text(s, 'HR App\nDemo Showcase', 0.72, 0.82, 4.7, 1.32, 35, WHITE, True, font='Aptos Display')
add_text(s, 'Employee Center → Time → Payroll → Benefit\nพร้อม approval path สำหรับ HR/Manager/HRBP/SPD', 0.76, 2.33, 4.75, 0.62, 15.5, RGBColor(201,230,226))
add_chip(s, 'SHOWCASE PACK', 0.78, 3.22, 1.65, AMBER, NAVY)
add_chip(s, 'FRONTEND DEMO', 2.62, 3.22, 1.75, TEAL, WHITE)
# Layered screenshots — keep cover readable at meeting-thumbnail size
add_image_card(s, imgs['quick'], 5.85, 0.72, 5.65, 3.95, 'Quick Approve', crop_bottom=0.20)
add_image_card(s, imgs['time'], 7.05, 4.55, 4.15, 1.88, 'Time', crop_bottom=0.04)
add_image_card(s, imgs['profile'], 10.62, 3.86, 1.65, 1.86, 'EC', crop_bottom=0.44)
# stats
for i,(num,lab) in enumerate([('4','modules'),('5','demo paths'),('0','backend claims')]):
    x=0.78+i*1.62
    add_round_rect(s, x, 4.18, 1.34, 0.96, color=RGBColor(20,48,72), line=RGBColor(37,74,99))
    add_text(s, num, x+0.14, 4.32, 0.45, 0.38, 24, MINT, True, 'left')
    add_text(s, lab, x+0.60, 4.42, 0.68, 0.24, 9.5, WHITE)
add_text(s, 'Prepared for app demo / design sign-off discussion', 0.78, 6.54, 4.2, 0.28, 11.5, RGBColor(198,219,226))
add_footer(s, 1)

# Slide 2 flow
s = prs.slides.add_slide(blank)
set_fill(s.background, CREAM)
add_title(s, 'Demo storyline', 'โชว์ให้เห็น end-to-end ไม่ใช่แค่หน้าจอเดี่ยว')
steps=[('1','Employee self-service','ดู profile, document, benefit, time request','profile'),('2','Manager action','ตรวจรายการและ approve / return','quick'),('3','HR / Payroll readiness','setup, review, caveat, control point','payroll')]
for i,(no,title,cap,img) in enumerate(steps):
    x=0.72+i*4.12
    add_round_rect(s,x,1.72,3.45,4.55,WHITE,BORDER)
    add_chip(s, f'STEP {no}', x+0.25,1.98,0.9, TEAL if i<2 else AMBER, WHITE if i<2 else NAVY)
    add_text(s,title,x+0.25,2.45,2.9,0.35,18,INK,True)
    add_text(s,cap,x+0.25,2.88,2.9,0.42,10.5,MUTED)
    add_image_card(s, imgs[img], x+0.25,3.48,2.95,2.08, None, crop_bottom=0.36 if img!='time' else 0.02)
    if i<2:
        add_text(s,'→',x+3.55,3.62,0.42,0.5,24,TEAL,True,'center')
add_footer(s,2)

# Slide 3 modules matrix
s = prs.slides.add_slide(blank)
set_fill(s.background, WHITE)
add_title(s, '4 modules ready for showcase', 'แต่ละ module มี route + mock path สำหรับ stakeholder review')
cards=[('Employee Center','Profile, documents, lifecycle copy',imgs['profile'],0.46),('Time','Attendance / timesheet / leave path',imgs['time'],0.02),('Payroll','Setup / processing / review cues',imgs['payroll'],0.46),('Benefit','Enrollment / claims / admin view',imgs['benefits'],0.42)]
positions=[(0.72,1.55),(6.95,1.55),(0.72,4.45),(6.95,4.45)]
for (title,cap,img,crop),(x,y) in zip(cards,positions):
    add_round_rect(s,x,y,5.65,2.25,CREAM,BORDER)
    add_image_card(s,img,x+0.18,y+0.18,2.38,1.88,None,crop_bottom=crop)
    add_text(s,title,x+2.8,y+0.32,2.5,0.32,18,INK,True)
    add_text(s,cap,x+2.8,y+0.78,2.45,0.42,10.8,MUTED)
    add_chip(s,'demo surface',x+2.8,y+1.42,1.18,TEAL,WHITE)
    add_chip(s,'sign-off ready',x+4.05,y+1.42,1.28,GREEN,WHITE)
add_footer(s,3)

# Slide 4 approval path
s = prs.slides.add_slide(blank)
set_fill(s.background, RGBColor(245,248,250))
add_title(s, 'Approval experience: one place to decide', 'Quick Approve เป็น demo anchor สำหรับ Manager / HRBP / SPD')
add_image_card(s, imgs['hrbp'], 0.72, 1.54, 7.15, 4.62, 'HRBP Quick Approve', crop_bottom=0.00)
for i,(role,task,color) in enumerate([('Manager','Approve / return employee request',TEAL),('HRBP','Bulk review + risk language',NAVY),('SPD','Policy / compliance checkpoint',AMBER)]):
    y=1.78+i*1.32
    add_round_rect(s,8.25,y,4.15,0.94,WHITE,BORDER)
    add_chip(s,role,8.48,y+0.24,0.92,color,WHITE if color!=AMBER else NAVY)
    add_text(s,task,9.55,y+0.25,2.45,0.3,12.5,INK,True)
add_text(s,'Design intent: ลด ambiguity ระหว่าง “เห็นข้อมูล”, “ตัดสินใจ”, และ “ส่งต่อ backend/workflow later”',8.35,5.92,3.78,0.6,12,MUTED)
add_footer(s,4)

# Slide 5 demo route
s = prs.slides.add_slide(blank)
set_fill(s.background, WHITE)
add_title(s, 'Suggested live demo route', 'ใช้ sequence นี้เพื่อให้ story ลื่นและไม่หลุด scope')
route=[('01','Home / Employee Center','ยืนยัน persona + employee data'),('02','Time / Leave','โชว์ request + policy caveat'),('03','Quick Approve','Manager action + return reason'),('04','Payroll','Setup/review view, no payroll correctness claim'),('05','Benefit','Benefit hub + admin/claim surface')]
for i,(no,title,cap) in enumerate(route):
    x=0.88+i*2.43
    add_round_rect(s,x,1.8,1.92,3.65,CREAM,BORDER)
    add_text(s,no,x+0.18,2.03,0.55,0.45,26,TEAL,True)
    add_text(s,title,x+0.18,2.72,1.45,0.58,14,INK,True)
    add_text(s,cap,x+0.18,3.55,1.48,0.62,10.2,MUTED)
    if i<4:
        add_text(s,'→',x+1.92,3.22,0.45,0.35,18,AMBER,True,'center')
add_image_card(s, imgs['benefits'], 1.22, 5.82, 3.0, 0.88, 'Benefit', crop_bottom=0.58)
add_image_card(s, imgs['quick'], 5.15, 5.82, 3.0, 0.88, 'Approve', crop_bottom=0.58)
add_image_card(s, imgs['payroll'], 9.10, 5.82, 3.0, 0.88, 'Payroll', crop_bottom=0.58)
add_footer(s,5)

# Slide 6 boundaries
s = prs.slides.add_slide(blank)
set_fill(s.background, CREAM)
add_title(s, 'What this demo proves / does not prove', 'ตั้ง expectation ให้ถูกก่อนเข้าห้อง demo')
left=[('Visual readiness','UI looks coherent enough for stakeholder walkthrough'),('Persona coverage','Employee, Manager, HRBP, SPD/Admin paths visible'),('Design sign-off pack','Screens + journey + caveats ready for review')]
right=[('No backend/database claim','Static/mock frontend only'),('No payroll correctness claim','Numbers and calculations are non-binding'),('No workflow integration claim','SF/workflow/audit are backend-later contracts')]
add_round_rect(s,0.82,1.65,5.65,4.65,WHITE,BORDER)
add_round_rect(s,6.86,1.65,5.65,4.65,WHITE,BORDER)
add_chip(s,'PROVES',1.08,1.95,1.0,TEAL,WHITE)
add_chip(s,'DOES NOT PROVE',7.12,1.95,1.5,AMBER,NAVY)
for i,(h,b) in enumerate(left):
    y=2.62+i*1.05
    add_text(s,'✓',1.1,y,0.26,0.25,17,GREEN,True)
    add_text(s,h,1.5,y,3.6,0.25,13.5,INK,True)
    add_text(s,b,1.5,y+0.31,4.25,0.25,10.4,MUTED)
for i,(h,b) in enumerate(right):
    y=2.62+i*1.05
    add_text(s,'—',7.16,y,0.26,0.25,17,AMBER,True)
    add_text(s,h,7.56,y,3.7,0.25,13.5,INK,True)
    add_text(s,b,7.56,y+0.31,4.35,0.25,10.4,MUTED)
add_footer(s,6)

# Slide 7 asks
s = prs.slides.add_slide(blank)
set_fill(s.background, DEEP)
add_title(s, 'Meeting focus', 'เป้าหมายคือ alignment และ design approval ก่อน backend', dark=True)
asks=[('1','Confirm demo route','ลำดับหน้าจอและ persona ที่จะใช้ในห้อง show case'),('2','Validate caveat wording','ข้อความ sensitive data / payroll / backend-later ต้องไม่ over-claim'),('3','Capture decisions','ประเด็นที่ต้องเข้ารอบ backend, integration, policy, data owner')]
for i,(no,h,b) in enumerate(asks):
    x=0.88+i*4.05
    add_round_rect(s,x,2.05,3.35,3.65,RGBColor(20,48,72),RGBColor(43,86,112))
    add_text(s,no,x+0.28,2.35,0.55,0.55,32,MINT,True)
    add_text(s,h,x+0.32,3.15,2.55,0.35,17,WHITE,True)
    add_text(s,b,x+0.32,3.72,2.55,0.68,11.5,RGBColor(201,225,229))
add_text(s,'Recommended output: signed-off storyboard + issue list for backend-later work',0.92,6.58,5.6,0.28,11.5,RGBColor(170,194,202))
add_footer(s,7)

prs.save(OUT)
print(OUT)
